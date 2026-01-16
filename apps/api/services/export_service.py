import io
from datetime import datetime
from docx import Document
from docx.shared import Pt, Mm
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from packages.core_domain import Generation

class ExportService:
    @staticmethod
    def generate_docx(generation: Generation, content: str) -> io.BytesIO:
        doc = Document()
        
        # Получаем настройки форматирования из generation.settings_payload
        formatting = generation.settings_payload.get('formatting', {})
        
        # --- Настройки страницы (поля) ---
        sections = doc.sections
        margins = formatting.get('margins', {})
        for section in sections:
            section.top_margin = Mm(margins.get('top', 20))
            section.bottom_margin = Mm(margins.get('bottom', 20))
            section.left_margin = Mm(margins.get('left', 30))
            section.right_margin = Mm(margins.get('right', 10))

        # --- Настройка стиля по умолчанию ---
        style = doc.styles['Normal']
        font = style.font
        font.name = formatting.get('fontFamily', 'Times New Roman')
        font.size = Pt(formatting.get('fontSize', 14))
        
        paragraph_format = style.paragraph_format
        line_spacing = formatting.get('lineSpacing', 1.5)
        if line_spacing == 1.0:
            paragraph_format.line_spacing_rule = WD_LINE_SPACING.SINGLE
        elif line_spacing == 1.5:
            paragraph_format.line_spacing_rule = WD_LINE_SPACING.ONE_POINT_FIVE
        elif line_spacing == 2.0:
            paragraph_format.line_spacing_rule = WD_LINE_SPACING.DOUBLE
        else:
            paragraph_format.line_spacing = line_spacing
            
        paragraph_indent_cm = formatting.get('paragraphIndent', 1.25)
        paragraph_format.first_line_indent = Mm(paragraph_indent_cm * 10) # см в мм
        
        alignment = formatting.get('alignment', 'justify')
        if alignment == 'justify':
            paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        elif alignment == 'left':
            paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT
        elif alignment == 'center':
            paragraph_format.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # --- Title Page (ГОСТ-титульник) ---
        # 1. Шапка (Министерство + ВУЗ)
        header_info = generation.input_payload.get("title_page", {})
        university = header_info.get("university", "ГОСУДАРСТВЕННОЕ ОБРАЗОВАТЕЛЬНОЕ УЧРЕЖДЕНИЕ")
        
        p = doc.add_paragraph('МИНИСТЕРСТВО НАУКИ И ВЫСШЕГО ОБРАЗОВАНИЯ РОССИЙСКОЙ ФЕДЕРАЦИИ')
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.first_line_indent = 0
        p.paragraph_format.line_spacing = 1.0 # В шапке одинарный интервал
        
        p = doc.add_paragraph(university.upper())
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.first_line_indent = 0
        
        doc.add_paragraph('\n' * 5)
        
        # 2. Тип работы и Тема
        work_type = (generation.work_type or 'РАБОТА').upper()
        p = doc.add_paragraph(f'{work_type}')
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.first_line_indent = 0
        if p.runs:
            run = p.runs[0]
            run.bold = True
            run.font.size = Pt(16)

        p = doc.add_paragraph('по теме:')
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.first_line_indent = 0

        p = doc.add_paragraph(generation.title or 'Без названия')
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.first_line_indent = 0
        if p.runs:
            run = p.runs[0]
            run.bold = True
            run.font.size = Pt(16)
        
        doc.add_paragraph('\n' * 8)
        
        # 3. Сведения об исполнителе
        full_name = header_info.get("full_name", "Студент И.И.")
        teacher_name = header_info.get("teacher_name", "Преподаватель П.П.")
        
        p = doc.add_paragraph(f'Выполнил: {full_name}')
        p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        p.paragraph_format.first_line_indent = 0
        
        p = doc.add_paragraph(f'Проверил: {teacher_name}')
        p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        p.paragraph_format.first_line_indent = 0
        
        doc.add_paragraph('\n' * 5)
        
        # 4. Город и Год
        city = header_info.get("city", "Москва")
        year = datetime.now().year
        p = doc.add_paragraph(f'{city}, {year}')
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.first_line_indent = 0
        
        doc.add_page_break()
        
        # --- Main Content ---
        for part in content.split('\n\n'):
            if part.strip() == '': continue
            
            if part.startswith('# '):
                h = doc.add_heading(part.replace('# ', ''), level=1)
                h.alignment = WD_ALIGN_PARAGRAPH.CENTER # Заголовки по центру
                h.paragraph_format.first_line_indent = 0
            elif part.startswith('## '):
                h = doc.add_heading(part.replace('## ', ''), level=2)
                h.alignment = WD_ALIGN_PARAGRAPH.LEFT
                h.paragraph_format.first_line_indent = Mm(12.5)
            else:
                p = doc.add_paragraph(part)
        
        file_stream = io.BytesIO()
        doc.save(file_stream)
        file_stream.seek(0)
        return file_stream

    @staticmethod
    def generate_pdf(generation: Generation, content: str) -> io.BytesIO:
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        
        pdf.cell(200, 10, txt="ЗАЧЁТ - РЕЗУЛЬТАТ ГЕНЕРАЦИИ", ln=True, align='C')
        pdf.ln(20)
        pdf.cell(200, 10, txt=f"Тема: {generation.title or 'Без названия'}", ln=True, align='C')
        pdf.ln(10)
        pdf.cell(200, 10, txt=f"Тип: {generation.module.value}", ln=True, align='C')
        pdf.ln(20)
        
        pdf.multi_cell(0, 10, txt=content)
        
        return io.BytesIO(pdf.output())

    @staticmethod
    def generate_pptx(generation: Generation, content: str) -> io.BytesIO:
        prs = Presentation()
        
        # --- Title Slide ---
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = generation.title or "Зачёт - Презентация"
        subtitle.text = f"Тип: {generation.module.value}\nСтиль: {generation.input_payload.get('presentation_style', 'academic')}"
        
        # --- Content Slides ---
        # Разбиваем контент на слайды по заголовкам ##
        sections = content.split('## ')
        for section in sections[1:]: # Пропускаем текст до первого заголовка
            lines = section.split('\n')
            slide_title_text = lines[0]
            slide_content_text = "\n".join(lines[1:]).strip()
            
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            slide.shapes.title.text = slide_title_text
            
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.text = slide_content_text
            
        file_stream = io.BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        return file_stream

export_service = ExportService()
